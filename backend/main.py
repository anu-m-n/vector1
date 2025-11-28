import sys
import os
import uuid
import re
import io
from typing import List, Optional, Dict, Any
from pypdf import PdfReader
from pptx import Presentation  # NEW: For PowerPoint
from docx import Document      # NEW: For Word Docs

# Fix imports
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

# Import Database Logic
from backend.database import vector_collection, conn
from backend.embeddings import get_embedding

app = FastAPI(title="Devfolio Hybrid DB + Uploads")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- HELPER: ROBUST FILE READER (NOW SUPPORTS PPTX & DOCX) ---
async def extract_text_safe(file: UploadFile) -> str:
    """
    Safely reads PDF, TXT, PPTX, and DOCX files.
    """
    try:
        # Read the file bytes into memory
        contents = await file.read()
        file_stream = io.BytesIO(contents)
        text = ""

        if file.filename.endswith(".pdf"):
            reader = PdfReader(file_stream)
            for page in reader.pages:
                if page.extract_text():
                    text += page.extract_text() + "\n"
            
        elif file.filename.endswith(".txt"):
            text = contents.decode("utf-8")

        elif file.filename.endswith(".pptx"):
            prs = Presentation(file_stream)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        elif file.filename.endswith(".docx"):
            doc = Document(file_stream)
            text = "\n".join([para.text for para in doc.paragraphs])
        
        else:
            print(f"Unsupported file type: {file.filename}")
            return "" 
            
        return text

    except Exception as e:
        print(f"ERROR reading file: {str(e)}")
        return ""

def extract_keywords(text):
    # Graph Logic: Find capitalized words (Entities)
    raw_words = re.findall(r'\b[A-Z][a-zA-Z]+\b', text)
    blacklist = {"The", "This", "That", "And", "For", "With", "But", "Redis", "Graph", "Slide", "Click", "Title"}
    return list(set([w for w in raw_words if w not in blacklist and len(w) > 3]))

# --- DATA MODELS ---
class NodeCreate(BaseModel):
    text: str
    metadata: Optional[Dict[str, Any]] = {}
    embedding: Optional[List[float]] = None

class HybridSearchRequest(BaseModel):
    query_text: str
    vector_weight: float = 0.5
    graph_weight: float = 0.5
    top_k: int = 5

# --- API ENDPOINTS ---

@app.post("/upload")
async def upload_document(file: UploadFile = File(...)):
    print(f"Received upload: {file.filename}")
    
    text = await extract_text_safe(file)
    
    if len(text.strip()) < 5:
        raise HTTPException(status_code=400, detail=f"File {file.filename} is empty or type not supported.")

    doc_id = str(uuid.uuid4())
    safe_title = file.filename.replace("'", "").replace('"', "")

    # A. Store in Vector DB
    vector = get_embedding(text)
    vector_collection.add(
        documents=[text[:2000]], 
        embeddings=[vector],
        ids=[doc_id],
        metadatas=[{"title": file.filename, "source": "User Upload"}]
    )

    # B. Store in Graph DB
    safe_text = text[:50].replace("'", "")
    conn.execute(f"MERGE (d:Document {{id: '{doc_id}', title: '{safe_title}', text: '{safe_text}'}})")

    # C. Auto-Link Entities
    keywords = extract_keywords(text[:5000])
    count = 0
    for word in keywords:
        word_id = word.lower()
        safe_word = word.replace("'", "")
        conn.execute(f"MERGE (e:Entity {{id: '{word_id}', name: '{safe_word}'}})")
        conn.execute(f"MATCH (d:Document), (e:Entity) WHERE d.id = '{doc_id}' AND e.id = '{word_id}' MERGE (d)-[:MENTIONS]->(e)")
        count += 1

    return {"status": "success", "id": doc_id, "entities_found": count, "message": "Processed successfully"}

# ... (Keep existing CRUD and SEARCH endpoints here) ...
@app.post("/nodes", status_code=201)
def create_node(item: NodeCreate):
    node_id = str(uuid.uuid4())
    vec = item.embedding if item.embedding else get_embedding(item.text)
    vector_collection.add(documents=[item.text], embeddings=[vec], ids=[node_id], metadatas=[item.metadata])
    safe_text = item.text.replace("'", "")[:50]
    conn.execute(f"MERGE (n:Document {{id: '{node_id}', title: '{safe_text}'}})")
    return {"id": node_id, "text": item.text}

@app.post("/search/hybrid")
def hybrid_search(req: HybridSearchRequest):
    query_vec = get_embedding(req.query_text)
    vec_res = vector_collection.query(query_embeddings=[query_vec], n_results=req.top_k * 2)
    results = []
    if vec_res['ids']:
        ids = vec_res['ids'][0]
        texts = vec_res['documents'][0]
        dists = vec_res['distances'][0]
        metas = vec_res['metadatas'][0]
        for i, doc_id in enumerate(ids):
            v_score = 1.0 / (1.0 + dists[i])
            g_score = 0.0
            entities = []
            try:
                q = f"MATCH (d:Document)-[:MENTIONS]->(e:Entity) WHERE d.id = '{doc_id}' RETURN e.name"
                kuzu_res = conn.execute(q)
                while kuzu_res.has_next(): entities.append(kuzu_res.get_next()[0])
                if len(entities) > 0: g_score = min(len(entities) * 0.1, 1.0)
            except: pass
            final = (req.vector_weight * v_score) + (req.graph_weight * g_score)
            results.append({"id": doc_id, "text": texts[i], "title": metas[i].get('title', 'Unknown'), "score": final, "details": f"Entities: {', '.join(entities[:3])}..."})
    results.sort(key=lambda x: x['score'], reverse=True)
    return {"results": results[:req.top_k]}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)